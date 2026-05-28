from flask import Flask, send_from_directory, request, send_file, session, redirect, url_for, render_template
from pptx import Presentation
from datetime import datetime
from PIL import Image
from functools import wraps
import os
import json
import zipfile
import uuid
import hmac

app = Flask(__name__, template_folder=".")

# IMPORTANTE:
# En Render configura SECRET_KEY como variable de entorno.
# Ejemplo: SECRET_KEY = una_clave_larga_y_privada
app.secret_key = os.environ.get("SECRET_KEY", "CAMBIAR_SECRET_KEY_EN_RENDER")

PLANTILLA = "plantillas/RESULTADOS_2025_PLANTILLA.pptx"
CARPETA_SALIDAS = "salidas"
CARPETA_UPLOADS = "uploads"
ARCHIVO_REGISTROS = "registros.json"
ARCHIVO_USUARIOS = "usuarios.json"

AREAS_OPERATIVAS = [
    "PEP_BJ", "PEP_PLAYA", "PEP_TUL", "PEP_COZ", "PEP_OPB", "PEP_CAMINOS",
    "POLICIA RURAL", "GPO_CENTURION", "GPO_ORION", "GPO_PDI", "GPO_JAGUAR",
    "MUN_FCP", "MUN_LC", "MUN_PLAYA", "MUN_TUL", "MUN_PM", "MUN_JMM",
    "MUN_IM", "MUN_COZ", "MUN_BJ", "MUN_RURAL", "MUN_BACALAR"
]


def nombre_variable_password(area):
    """Convierte PEP_BJ -> PASS_PEP_BJ y POLICIA RURAL -> PASS_POLICIA_RURAL."""
    limpio = str(area).upper().replace(" ", "_").replace("-", "_")
    return f"PASS_{limpio}"


def cargar_usuarios():
    """
    Carga usuarios por dos vías:
    1) usuarios.json, útil para pruebas locales.
    2) Variables de entorno PASS_..., recomendado para Render.

    Formato usuarios.json:
    {
      "ESTADISTICAS": {"password":"clave", "area":"ESTADISTICAS", "rol":"ADMIN"},
      "PEP_BJ": {"password":"clave", "area":"PEP_BJ", "rol":"AREA"}
    }
    """
    usuarios = {}

    if os.path.exists(ARCHIVO_USUARIOS):
        try:
            with open(ARCHIVO_USUARIOS, "r", encoding="utf-8") as f:
                data = json.load(f)

            for usuario, info in data.items():
                if isinstance(info, str):
                    area = usuario
                    rol = "ADMIN" if usuario == "ESTADISTICAS" else "AREA"
                    password = info
                else:
                    area = info.get("area", usuario)
                    rol = info.get("rol", "AREA")
                    password = info.get("password", "")

                if password:
                    usuarios[usuario] = {
                        "password": str(password),
                        "area": str(area),
                        "rol": str(rol).upper()
                    }
        except Exception as err:
            print("Error leyendo usuarios.json:", err)

    # Variables de entorno. Si existen, tienen prioridad sobre usuarios.json.
    pass_admin = os.environ.get("PASS_ESTADISTICAS")
    if pass_admin:
        usuarios["ESTADISTICAS"] = {
            "password": pass_admin,
            "area": "ESTADISTICAS",
            "rol": "ADMIN"
        }

    for area in AREAS_OPERATIVAS:
        password = os.environ.get(nombre_variable_password(area))
        if password:
            usuarios[area] = {
                "password": password,
                "area": area,
                "rol": "AREA"
            }

    return usuarios


def usuario_actual():
    return {
        "usuario": session.get("usuario", ""),
        "area": session.get("area", ""),
        "rol": session.get("rol", "")
    }


def es_admin():
    return session.get("rol") == "ADMIN"


def login_requerido(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not session.get("usuario"):
            return redirect(url_for("login"))
        return func(*args, **kwargs)
    return wrapper


def api_login_requerido(func):
    @wraps(func)
    def wrapper(*args, **kwargs):
        if not session.get("usuario"):
            return {"ok": False, "error": "Sesión no iniciada"}, 401
        return func(*args, **kwargs)
    return wrapper


@app.route("/login", methods=["GET", "POST"])
def login():
    usuarios = cargar_usuarios()
    error = ""

    if request.method == "POST":
        usuario = request.form.get("usuario", "").strip()
        password = request.form.get("password", "")

        info = usuarios.get(usuario)

        if info and hmac.compare_digest(password, info.get("password", "")):
            session.clear()
            session["usuario"] = usuario
            session["area"] = info.get("area", usuario)
            session["rol"] = info.get("rol", "AREA")
            return redirect(url_for("inicio"))

        error = "Usuario o contraseña incorrectos."

    if not usuarios and not error:
        error = "No hay usuarios configurados. Crea usuarios.json o variables PASS_... en Render."

    nombres = sorted(usuarios.keys(), key=lambda x: (x != "ESTADISTICAS", x))
    return render_template("login.html", usuarios=nombres, error=error)


@app.route("/logout")
def logout():
    session.clear()
    return redirect(url_for("login"))


@app.route("/")
@login_requerido
def inicio():
    return send_from_directory(".", "Detenidos.html")


@app.route("/<path:archivo>")
@login_requerido
def archivos(archivo):
    """
    Sirve recursos públicos del formulario solo si hay sesión iniciada.
    Bloquea archivos sensibles para que no se descarguen directo por URL.
    """
    archivo_norm = archivo.replace("\\", "/")

    bloqueados_exactos = {
        ARCHIVO_REGISTROS,
        ARCHIVO_USUARIOS,
        "usuarios_ejemplo.json",
        "app.py",
        "app_mejora_accesos.py"
    }

    carpetas_bloqueadas = (
        "uploads/",
        "salidas/",
        "plantillas/"
    )

    if archivo_norm in bloqueados_exactos or archivo_norm.startswith(carpetas_bloqueadas):
        return "No autorizado", 403

    if archivo_norm == "Detenidos.html":
        return send_from_directory(".", "Detenidos.html")

    return send_from_directory(".", archivo_norm)


def limpiar_nombre(texto):
    texto = str(texto or "SIN_FOLIO")
    for c in ['/', '\\', ':', '*', '?', '"', '<', '>', '|']:
        texto = texto.replace(c, "_")
    return texto.strip()


def formato_fecha_mx(fecha):
    if not fecha:
        return ""

    partes = str(fecha).split("-")

    if len(partes) == 3:
        return f"{partes[2]}/{partes[1]}/{partes[0]}"

    return str(fecha)


def construir_reemplazos(data):
    detenido = (
        f"{data.get('nombre', '')} "
        f"{data.get('ap_paterno', '')} "
        f"{data.get('ap_materno', '')}"
    ).strip()

    return {
        "<FECHA>": formato_fecha_mx(data.get("fecha", "")),
        "<DELITO>": data.get("hecho", ""),
        "<DETENIDO>": detenido,
        "<EDAD>": data.get("edad", ""),
        "<ALIAS>": data.get("alias", ""),

        "<DETENIDO1>": data.get("detenido1", ""),
        "<ALIAS1>": data.get("alias1", ""),
        "<DETENIDO2>": data.get("detenido2", ""),
        "<ALIAS2>": data.get("alias2", ""),
        "<DETENIDO3>": data.get("detenido3", ""),
        "<ALIAS3>": data.get("alias3", ""),
        "<DETENIDO4>": data.get("detenido4", ""),
        "<ALIAS4>": data.get("alias4", ""),

        "<ASEGURAMIENTO1>": data.get("aseguramiento_1", ""),
        "<ASEGURAMIENTO2>": data.get("aseguramiento_2", ""),
        "<ASEGURAMIENTO3>": data.get("aseguramiento_3", ""),
        "<ASEGURAMIENTO4>": data.get("aseguramiento_4", ""),
        "<ASEGURAMIENTO5>": data.get("aseguramiento_5", ""),

        "<INFORMACIONRELEVANTE>": data.get("observaciones_rnd", ""),
        "<CORP>": data.get("id_operativo", "")
    }


def reemplazar_texto_pptx(prs, reemplazos):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for clave, valor in reemplazos.items():
                        if clave in run.text:
                            run.text = run.text.replace(clave, str(valor or ""))


def borrar_shape(shape):
    elemento = shape._element
    elemento.getparent().remove(elemento)


def validar_area_registro(registro):
    """
    Regla central de seguridad:
    - ESTADISTICAS puede trabajar con todos los registros.
    - Un área normal solo puede trabajar con registros de su misma área.
    """
    if es_admin():
        return True

    area_sesion = session.get("area", "")
    area_registro = str(registro.get("id_operativo", "")).strip()

    return bool(area_sesion and area_registro == area_sesion)


def preparar_registro_para_guardar(registro):
    """
    Valida o completa metadatos antes de guardar.
    Para usuarios AREA se exige que id_operativo coincida con su sesión.
    """
    if not isinstance(registro, dict):
        return None, ("Registro inválido", 400)

    if es_admin():
        registro["_area_captura"] = registro.get("id_operativo", "")
    else:
        area_sesion = session.get("area", "")
        area_formulario = str(registro.get("id_operativo", "")).strip()

        if area_formulario != area_sesion:
            return None, ("No autorizado: el área del registro no coincide con tu sesión", 403)

        registro["_area_captura"] = area_sesion

    registro["_usuario_captura"] = session.get("usuario", "")
    return registro, None


def registros_visibles(registros):
    if es_admin():
        return registros

    area_sesion = session.get("area", "")
    return [
        r for r in registros
        if str(r.get("id_operativo", "")).strip() == area_sesion
        or str(r.get("_area_captura", "")).strip() == area_sesion
    ]


def registros_visibles_con_indices(registros):
    if es_admin():
        return list(enumerate(registros))

    area_sesion = session.get("area", "")
    return [
        (i, r) for i, r in enumerate(registros)
        if str(r.get("id_operativo", "")).strip() == area_sesion
        or str(r.get("_area_captura", "")).strip() == area_sesion
    ]


def guardar_imagenes_subidas():
    os.makedirs(CARPETA_UPLOADS, exist_ok=True)

    imagenes = {}

    for i in range(1, 5):
        archivo = request.files.get(f"img{i}")

        if archivo and archivo.filename:
            ext = os.path.splitext(archivo.filename)[1].lower()

            if ext in [".jpg", ".jpeg", ".png"]:
                nombre = f"img{i}_{uuid.uuid4().hex}{ext}"
                ruta = os.path.join(CARPETA_UPLOADS, nombre)
                archivo.save(ruta)
                imagenes[f"<IMAGEN_{i}>"] = ruta.replace("\\", "/")

    if not imagenes:
        archivos = request.files.getlist("imagenes[]")

        for i, archivo in enumerate(archivos[:4], start=1):
            if archivo and archivo.filename:
                ext = os.path.splitext(archivo.filename)[1].lower()

                if ext in [".jpg", ".jpeg", ".png"]:
                    nombre = f"img{i}_{uuid.uuid4().hex}{ext}"
                    ruta = os.path.join(CARPETA_UPLOADS, nombre)
                    archivo.save(ruta)
                    imagenes[f"<IMAGEN_{i}>"] = ruta.replace("\\", "/")

    print("IMAGENES RECIBIDAS:", imagenes)

    return imagenes


@app.route("/guardar_imagenes_registro", methods=["POST"])
@api_login_requerido
def guardar_imagenes_registro():
    os.makedirs(CARPETA_UPLOADS, exist_ok=True)

    rutas = {}

    for i in range(1, 5):
        archivo = request.files.get(f"img{i}")

        if archivo and archivo.filename:
            ext = os.path.splitext(archivo.filename)[1].lower()

            if ext in [".jpg", ".jpeg", ".png"]:
                nombre = f"registro_img{i}_{uuid.uuid4().hex}{ext}"
                ruta = os.path.join(CARPETA_UPLOADS, nombre)
                archivo.save(ruta)

                rutas[f"img{i}_ruta"] = ruta.replace("\\", "/")

    print("RUTAS GUARDADAS REGISTRO:", rutas)

    return rutas


def imagenes_desde_registro(registro):
    imagenes = {}

    for i in range(1, 5):
        ruta = registro.get(f"img{i}_ruta", "")

        if ruta and os.path.exists(ruta):
            imagenes[f"<IMAGEN_{i}>"] = ruta

    print("IMAGENES DESDE REGISTRO:", imagenes)

    return imagenes


def insertar_imagenes_pptx(prs, imagenes):
    for slide in prs.slides:
        shapes = list(slide.shapes)

        for shape in shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            texto_shape = shape.text or ""

            for marcador, ruta_img in imagenes.items():
                if marcador in texto_shape:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    borrar_shape(shape)

                    img = Image.open(ruta_img)
                    img_w, img_h = img.size

                    ratio_img = img_w / img_h
                    ratio_box = width / height

                    if ratio_img > ratio_box:
                        nuevo_width = width
                        nuevo_height = int(width / ratio_img)
                    else:
                        nuevo_height = height
                        nuevo_width = int(height * ratio_img)

                    nuevo_left = left + int((width - nuevo_width) / 2)
                    nuevo_top = top + int((height - nuevo_height) / 2)

                    slide.shapes.add_picture(
                        ruta_img,
                        nuevo_left,
                        nuevo_top,
                        width=nuevo_width,
                        height=nuevo_height
                    )

                    break


def crear_pptx_desde_plantilla(data, ruta_salida, imagenes=None):
    prs = Presentation(PLANTILLA)

    if imagenes:
        insertar_imagenes_pptx(prs, imagenes)

    reemplazos = construir_reemplazos(data)
    reemplazar_texto_pptx(prs, reemplazos)

    prs.save(ruta_salida)


@app.route("/generar_pptx", methods=["POST"])
@api_login_requerido
def generar_pptx():
    os.makedirs(CARPETA_SALIDAS, exist_ok=True)

    data = request.form.to_dict()

    if not validar_area_registro(data):
        return "No autorizado para generar PPTX de otra área", 403

    imagenes = guardar_imagenes_subidas()

    folio = limpiar_nombre(data.get("folio_iph", "SIN_FOLIO"))
    nombre_archivo = f"{folio}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    ruta_salida = os.path.join(CARPETA_SALIDAS, nombre_archivo)

    crear_pptx_desde_plantilla(data, ruta_salida, imagenes)

    return send_file(ruta_salida, as_attachment=True)


@app.route("/generar_pptx_zip", methods=["POST"])
@api_login_requerido
def generar_pptx_zip():
    # Seguridad: no confiamos en records_json del navegador.
    # El ZIP se genera desde los registros guardados en servidor y filtrados por sesión.
    registros = registros_visibles(leer_registros_servidor())

    if not registros:
        return "No hay registros autorizados para generar ZIP", 400

    os.makedirs(CARPETA_SALIDAS, exist_ok=True)

    nombre_zip = f"PPTX_TODOS_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    ruta_zip = os.path.join(CARPETA_SALIDAS, nombre_zip)

    with zipfile.ZipFile(ruta_zip, "w", zipfile.ZIP_DEFLATED) as zipf:
        for i, registro in enumerate(registros, start=1):
            folio = limpiar_nombre(registro.get("folio_iph", f"REGISTRO_{i}"))

            nombre_pptx = f"{i:03d}_{folio}.pptx"
            ruta_pptx = os.path.join(CARPETA_SALIDAS, nombre_pptx)

            imagenes_registro = imagenes_desde_registro(registro)

            crear_pptx_desde_plantilla(
                registro,
                ruta_pptx,
                imagenes_registro
            )

            zipf.write(ruta_pptx, nombre_pptx)

    return send_file(ruta_zip, as_attachment=True)


def leer_registros_servidor():
    if not os.path.exists(ARCHIVO_REGISTROS):
        return []

    try:
        with open(ARCHIVO_REGISTROS, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception as err:
        print("Error leyendo registros:", err)
        return []


def guardar_registros_servidor(registros):
    with open(ARCHIVO_REGISTROS, "w", encoding="utf-8") as f:
        json.dump(registros, f, ensure_ascii=False, indent=2)


@app.route("/api/sesion", methods=["GET"])
@api_login_requerido
def api_sesion():
    return usuario_actual()


@app.route("/api/registros", methods=["GET"])
@api_login_requerido
def api_listar_registros():
    registros = leer_registros_servidor()
    return registros_visibles(registros)


@app.route("/api/registros", methods=["POST"])
@api_login_requerido
def api_guardar_registro():
    registro = request.get_json()

    if not registro:
        return {"ok": False, "error": "Registro vacío"}, 400

    registro, error = preparar_registro_para_guardar(registro)
    if error:
        mensaje, status = error
        return {"ok": False, "error": mensaje}, status

    registros = leer_registros_servidor()
    registros.append(registro)
    guardar_registros_servidor(registros)

    return {"ok": True, "total": len(registros_visibles(registros))}


@app.route("/api/registros", methods=["DELETE"])
@api_login_requerido
def api_borrar_registros():
    registros = leer_registros_servidor()

    if es_admin():
        guardar_registros_servidor([])
        return {"ok": True, "total": 0}

    area_sesion = session.get("area", "")
    registros_filtrados = [
        r for r in registros
        if not (
            str(r.get("id_operativo", "")).strip() == area_sesion
            or str(r.get("_area_captura", "")).strip() == area_sesion
        )
    ]

    guardar_registros_servidor(registros_filtrados)
    return {"ok": True, "total": 0}


@app.route("/api/registros/<int:index>", methods=["DELETE"])
@api_login_requerido
def api_eliminar_registro(index):
    registros = leer_registros_servidor()
    visibles = registros_visibles_con_indices(registros)

    if index < 0 or index >= len(visibles):
        return {"ok": False, "error": "Índice inválido"}, 404

    indice_real = visibles[index][0]
    registros.pop(indice_real)
    guardar_registros_servidor(registros)

    return {"ok": True, "total": len(registros_visibles(registros))}


if __name__ == "__main__":
    app.run(
        host="0.0.0.0",
        port=int(os.environ.get("PORT", 5000)),
        debug=os.environ.get("FLASK_DEBUG", "0") == "1"
    )
