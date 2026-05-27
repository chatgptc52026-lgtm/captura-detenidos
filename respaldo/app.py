from flask import Flask, send_from_directory, request, send_file
from pptx import Presentation
from datetime import datetime
from PIL import Image
import os
import json
import zipfile
import uuid

app = Flask(__name__)

PLANTILLA = "plantillas/RESULTADOS_2025_PLANTILLA.pptx"
CARPETA_SALIDAS = "salidas"
CARPETA_UPLOADS = "uploads"
ARCHIVO_REGISTROS = "registros.json"


@app.route("/")
def inicio():
    return send_from_directory(".", "Detenidos.html")


@app.route("/<path:archivo>")
def archivos(archivo):
    return send_from_directory(".", archivo)


def limpiar_nombre(texto):
    texto = str(texto or "SIN_FOLIO")
    for c in ['/', '\\', ':', '*', '?', '"', '<', '>', '|']:
        texto = texto.replace(c, "_")
    return texto.strip()


def construir_reemplazos(data):
    detenido = (
        f"{data.get('nombre', '')} "
        f"{data.get('ap_paterno', '')} "
        f"{data.get('ap_materno', '')}"
    ).strip()

    return {
        "<FECHA>": data.get("fecha", ""),
        "<DELITO>": data.get("hecho", ""),
        "<DETENIDO>": detenido,
        "<EDAD>": data.get("edad", ""),
        "<ALIAS>": data.get("alias", ""),

        "<DETENIDO1>": data.get("detenido1", ""),
        "<ALIAS1>": data.get("alias1", ""),
        "<DETENIDO2>": data.get("detenido2", ""),
        "<ALIAS2>": data.get("alias2", ""),
        "<DETENIDO3>": data.get("detenido3", ""),
        "<ALIAS3>": data.get("alias3", ""),
        "<DETENIDO4>": data.get("detenido4", ""),
        "<ALIAS4>": data.get("alias4", ""),

        "<ASEGURAMIENTO1>": data.get("aseguramiento_1", ""),
        "<ASEGURAMIENTO2>": data.get("aseguramiento_2", ""),
        "<ASEGURAMIENTO3>": data.get("aseguramiento_3", ""),
        "<ASEGURAMIENTO4>": data.get("aseguramiento_4", ""),
        "<ASEGURAMIENTO5>": data.get("aseguramiento_5", ""),

        "<INFORMACIONRELEVANTE>": data.get("observaciones_rnd", ""),
        "<CORP>": data.get("id_operativo", "")
    }


def reemplazar_texto_pptx(prs, reemplazos):
    for slide in prs.slides:
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    for clave, valor in reemplazos.items():
                        if clave in run.text:
                            run.text = run.text.replace(clave, str(valor or ""))


def borrar_shape(shape):
    elemento = shape._element
    elemento.getparent().remove(elemento)


def guardar_imagenes_subidas():
    os.makedirs(CARPETA_UPLOADS, exist_ok=True)

    imagenes = {}

    for i in range(1, 5):
        archivo = request.files.get(f"img{i}")

        if archivo and archivo.filename:
            ext = os.path.splitext(archivo.filename)[1].lower()

            if ext in [".jpg", ".jpeg", ".png"]:
                nombre = f"img{i}_{uuid.uuid4().hex}{ext}"
                ruta = os.path.join(CARPETA_UPLOADS, nombre)
                archivo.save(ruta)
                imagenes[f"<IMAGEN_{i}>"] = ruta.replace("\\", "/")

    if not imagenes:
        archivos = request.files.getlist("imagenes[]")

        for i, archivo in enumerate(archivos[:4], start=1):
            if archivo and archivo.filename:
                ext = os.path.splitext(archivo.filename)[1].lower()

                if ext in [".jpg", ".jpeg", ".png"]:
                    nombre = f"img{i}_{uuid.uuid4().hex}{ext}"
                    ruta = os.path.join(CARPETA_UPLOADS, nombre)
                    archivo.save(ruta)
                    imagenes[f"<IMAGEN_{i}>"] = ruta.replace("\\", "/")

    print("IMAGENES RECIBIDAS:", imagenes)

    return imagenes


@app.route("/guardar_imagenes_registro", methods=["POST"])
def guardar_imagenes_registro():
    os.makedirs(CARPETA_UPLOADS, exist_ok=True)

    rutas = {}

    for i in range(1, 5):
        archivo = request.files.get(f"img{i}")

        if archivo and archivo.filename:
            ext = os.path.splitext(archivo.filename)[1].lower()

            if ext in [".jpg", ".jpeg", ".png"]:
                nombre = f"registro_img{i}_{uuid.uuid4().hex}{ext}"
                ruta = os.path.join(CARPETA_UPLOADS, nombre)
                archivo.save(ruta)

                rutas[f"img{i}_ruta"] = ruta.replace("\\", "/")

    print("RUTAS GUARDADAS REGISTRO:", rutas)

    return rutas


def imagenes_desde_registro(registro):
    imagenes = {}

    for i in range(1, 5):
        ruta = registro.get(f"img{i}_ruta", "")

        if ruta and os.path.exists(ruta):
            imagenes[f"<IMAGEN_{i}>"] = ruta

    print("IMAGENES DESDE REGISTRO:", imagenes)

    return imagenes


def insertar_imagenes_pptx(prs, imagenes):
    for slide in prs.slides:
        shapes = list(slide.shapes)

        for shape in shapes:
            if not getattr(shape, "has_text_frame", False):
                continue

            texto_shape = shape.text or ""

            for marcador, ruta_img in imagenes.items():
                if marcador in texto_shape:
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height

                    borrar_shape(shape)

                    img = Image.open(ruta_img)
                    img_w, img_h = img.size

                    ratio_img = img_w / img_h
                    ratio_box = width / height

                    if ratio_img > ratio_box:
                        nuevo_width = width
                        nuevo_height = int(width / ratio_img)
                    else:
                        nuevo_height = height
                        nuevo_width = int(height * ratio_img)

                    nuevo_left = left + int((width - nuevo_width) / 2)
                    nuevo_top = top + int((height - nuevo_height) / 2)

                    slide.shapes.add_picture(
                        ruta_img,
                        nuevo_left,
                        nuevo_top,
                        width=nuevo_width,
                        height=nuevo_height
                    )

                    break


def crear_pptx_desde_plantilla(data, ruta_salida, imagenes=None):
    prs = Presentation(PLANTILLA)

    if imagenes:
        insertar_imagenes_pptx(prs, imagenes)

    reemplazos = construir_reemplazos(data)
    reemplazar_texto_pptx(prs, reemplazos)

    prs.save(ruta_salida)


@app.route("/generar_pptx", methods=["POST"])
def generar_pptx():
    os.makedirs(CARPETA_SALIDAS, exist_ok=True)

    imagenes = guardar_imagenes_subidas()

    folio = limpiar_nombre(request.form.get("folio_iph", "SIN_FOLIO"))
    nombre_archivo = f"{folio}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
    ruta_salida = os.path.join(CARPETA_SALIDAS, nombre_archivo)

    crear_pptx_desde_plantilla(request.form, ruta_salida, imagenes)

    return send_file(ruta_salida, as_attachment=True)


@app.route("/generar_pptx_zip", methods=["POST"])
def generar_pptx_zip():
    records_json = request.form.get("records_json", "[]")
    registros = json.loads(records_json)

    if not registros:
        return "No hay registros para generar ZIP", 400

    os.makedirs(CARPETA_SALIDAS, exist_ok=True)

    nombre_zip = f"PPTX_TODOS_{datetime.now().strftime('%Y%m%d_%H%M%S')}.zip"
    ruta_zip = os.path.join(CARPETA_SALIDAS, nombre_zip)

    with zipfile.ZipFile(ruta_zip, "w", zipfile.ZIP_DEFLATED) as zipf:
        for i, registro in enumerate(registros, start=1):
            folio = limpiar_nombre(registro.get("folio_iph", f"REGISTRO_{i}"))

            nombre_pptx = f"{i:03d}_{folio}.pptx"
            ruta_pptx = os.path.join(CARPETA_SALIDAS, nombre_pptx)

            imagenes_registro = imagenes_desde_registro(registro)

            crear_pptx_desde_plantilla(
                registro,
                ruta_pptx,
                imagenes_registro
            )

            zipf.write(ruta_pptx, nombre_pptx)

    return send_file(ruta_zip, as_attachment=True)

def leer_registros_servidor():
    if not os.path.exists(ARCHIVO_REGISTROS):
        return []

    try:
        with open(ARCHIVO_REGISTROS, "r", encoding="utf-8") as f:
            return json.load(f)
    except:
        return []


def guardar_registros_servidor(registros):
    with open(ARCHIVO_REGISTROS, "w", encoding="utf-8") as f:
        json.dump(registros, f, ensure_ascii=False, indent=2)


@app.route("/api/registros", methods=["GET"])
def api_listar_registros():
    return leer_registros_servidor()


@app.route("/api/registros", methods=["POST"])
def api_guardar_registro():
    registro = request.get_json()

    if not registro:
        return {"ok": False, "error": "Registro vacío"}, 400

    registros = leer_registros_servidor()
    registros.append(registro)
    guardar_registros_servidor(registros)

    return {"ok": True, "total": len(registros)}


@app.route("/api/registros", methods=["DELETE"])
def api_borrar_registros():
    guardar_registros_servidor([])
    return {"ok": True}

@app.route("/api/registros/<int:index>", methods=["DELETE"])
def api_eliminar_registro(index):
    registros = leer_registros_servidor()

    if index < 0 or index >= len(registros):
        return {"ok": False, "error": "Índice inválido"}, 404

    registros.pop(index)
    guardar_registros_servidor(registros)

    return {"ok": True, "total": len(registros)}

if __name__ == "__main__":
    app.run(
    host="0.0.0.0",
    port=5000,
    debug=True
)